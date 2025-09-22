from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
import ebooklib
from bs4 import BeautifulSoup
from pdf2image import convert_from_path
import easyocr
import tempfile
import requests
from dotenv import load_dotenv
import os

load_dotenv()

OCR_reader = easyocr.Reader(['en'], gpu=False)

OCR_space_api = os.getenv("OCR_SPACE_API")

language='eng'

# ---- Extractors for each file type ---- #

# Initialize EasyOCR once (to avoid reloading model each call)

def extract_image(file_path: str) -> str:

    """Extract text from an image file using OCR."""

    payload = {
        'isOverlayRequired': False,
        'apikey': OCR_space_api,
        'language': language,
    }
    with open(file_path, 'rb') as f:
        response = requests.post(
            'https://api.ocr.space/parse/image',
            files={file_path: f},
            data=payload,
        )
    result = response.json()
    try:
        return result['ParsedResults'][0]['ParsedText']
    except Exception as e:
        print("Error:", e, result)
        return ""

def extract_txt(file_path: str) -> str:
    """Extract text from a plain text file."""
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text + '\n\n'


def extract_pdf(file_path: str) -> str:
    """
    Extract text from a PDF file.
    - Uses PdfReader for selectable text.
    - Falls back to OCR via pytesseract for scanned/image pages.
    """
    full_text = []
    try:
        reader = PdfReader(file_path)
        for index, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():  
                # If text is found on the page
                full_text.append(text.strip())
            else:
                # If page has no text, fallback to OCR
                images = convert_from_path(file_path, first_page=index + 1, last_page=index + 1)
                if images:
                    # Save page as a temporary image file
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=True) as tmp_img:
                        images[0].save(tmp_img.name, format="PNG")

                        # Send to OCR.space API
                        ocr_text = extract_image(tmp_img.name).strip()

                        if ocr_text:
                            full_text.append(ocr_text)
        return '\n'.join(full_text) + '\n\n'
    except FileNotFoundError:
        print(f"File not found at path '{file_path}'")
        return ''


def extract_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""
    try:
        doc = Document(file_path)
        full_text = [para.text for para in doc.paragraphs]
        return '\n'.join(full_text) + '\n\n'
    except Exception as e:
        print(f"Error processing DOCX file: {e}")
        return ''


def extract_excel(file_path: str) -> str:
    """
    Extract tabular content from CSV/Excel files.
    
    Args:
        file_path (str): Path to the file (.csv, .xls, .xlsx).
    """
    file_type = file_path.rsplit('.')[-1].lower()
    if file_type == 'csv':
        df = pd.read_csv(file_path)
    elif file_type in ['xls', 'xlsx', 'excel']:
        df = pd.read_excel(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
    return df.to_string() + '\n\n'


def extract_pptx(file_path: str) -> str:
    """Extract text from all slides in a PPTX presentation."""
    text = []
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text.append(shape.text)
        text.append('\n')  # Slide separator

    return '\n'.join(text) + '\n\n'


def extract_epub(file_path: str) -> str:
    """Extract text from an EPUB file."""
    text = []
    book = ebooklib.epub.read_epub(file_path)

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_body_content(), "html.parser")
            text.append(soup.get_text())

    return '\n'.join(text) + '\n\n'

# Mapping of supported file types to their extractor functions
EXTRACTORS = {
    "pdf": extract_pdf,
    "txt": extract_txt,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "epub": extract_epub,
    "xls": extract_excel,
    "xlsx": extract_excel,
    "csv": extract_excel,
    "jpg": extract_image,
    "jpeg": extract_image,
    "png": extract_image,
    "bmp": extract_image
}
